# -*- coding: utf-8 -*-
"""Generate the defense presentation for ML_Practice.

The deck is intentionally generated from code so metrics, figures, and layout
can be updated consistently after experiments or reports change.
"""
from pathlib import Path
import json

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentation" / "人工智能综合实训II_答辩汇报.pptx"

T1 = ROOT / "task1_ml_wine" / "outputs"
T2 = ROOT / "task2_alexnet_fmnist" / "outputs"
REPORT = ROOT / "report"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

GREEN = RGBColor(0x0B, 0x6B, 0x4D)
GREEN_DARK = RGBColor(0x06, 0x45, 0x36)
BLUE = RGBColor(0x25, 0x67, 0xA8)
AMBER = RGBColor(0xB8, 0x6B, 0x00)
RED = RGBColor(0xB8, 0x3A, 0x3A)
INK = RGBColor(0x1F, 0x2A, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF6, 0xF8, 0xFA)
LINE = RGBColor(0xD9, 0xDF, 0xE7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_CN = "Microsoft YaHei"
FONT_EN = "Aptos"


def safe_json(path, default=None):
    if not path.exists():
        return default
    with open(path, encoding="utf-8") as f:
        return json.load(f)


def set_run(run, size=18, color=INK, bold=False, font=FONT_CN):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.color.rgb = LINE
    shape.line.width = Pt(0.7)


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold)
    return box


def add_title(slide, title, subtitle=None, tag=None):
    add_text(slide, title, Inches(0.55), Inches(0.30), Inches(10.5), Inches(0.45),
             size=25, color=INK, bold=True)
    if subtitle:
        add_text(slide, subtitle, Inches(0.57), Inches(0.79), Inches(10.8), Inches(0.30),
                 size=10.5, color=MUTED)
    if tag:
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.25), Inches(0.33),
                                      Inches(1.5), Inches(0.35))
        set_fill(pill, GREEN, transparency=0)
        pill.line.color.rgb = GREEN
        add_text(slide, tag, Inches(11.33), Inches(0.39), Inches(1.35), Inches(0.20),
                 size=9.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.14),
                                  Inches(12.22), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.color.rgb = LINE


def add_footer(slide, idx):
    add_text(slide, "人工智能综合实训 II · 机器学习项目答辩", Inches(0.55), Inches(7.08),
             Inches(5.8), Inches(0.22), size=8.5, color=MUTED)
    add_text(slide, f"{idx:02d}", Inches(12.25), Inches(7.05), Inches(0.5), Inches(0.25),
             size=9.5, color=MUTED, align=PP_ALIGN.RIGHT)


def add_bullets(slide, items, x, y, w, h, size=15, color=INK, gap=4):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(gap)
        p.font.name = FONT_CN
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_card(slide, x, y, w, h, title, body=None, accent=GREEN, title_size=13,
             body_size=18, body_bold=True):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(card, WHITE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    add_text(slide, title, x + Inches(0.18), y + Inches(0.12), w - Inches(0.3), Inches(0.25),
             size=title_size, color=MUTED, bold=True)
    if body:
        add_text(slide, body, x + Inches(0.18), y + Inches(0.45), w - Inches(0.28), h - Inches(0.55),
                 size=body_size, color=INK, bold=body_bold, valign=MSO_ANCHOR.MIDDLE)
    return card


def add_image_fit(slide, path, x, y, w, h, border=True):
    path = Path(path)
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_fill(frame, WHITE)
    if border:
        frame.line.color.rgb = LINE
    else:
        frame.line.fill.background()
    if not path.exists():
        add_text(slide, f"缺少图片：{path.name}", x + Inches(0.2), y + h / 2 - Inches(0.15),
                 w - Inches(0.4), Inches(0.3), size=12, color=RED, align=PP_ALIGN.CENTER)
        return frame
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw = int(iw * scale)
    ph = int(ih * scale)
    px = x + (w - pw) / 2
    py = y + (h - ph) / 2
    pic = slide.shapes.add_picture(str(path), px, py, width=pw, height=ph)
    return pic


def add_table(slide, x, y, w, h, headers, rows, font_size=10.5):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h)
    table = table_shape.table
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN_DARK
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                set_run(run, size=font_size, color=WHITE, bold=True)
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC) if r % 2 else WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    set_run(run, size=font_size, color=INK, bold=False)
    for col in range(len(headers)):
        table.cell(0, col).margin_left = Inches(0.03)
        table.cell(0, col).margin_right = Inches(0.03)
    return table_shape


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xFB, 0xFC, 0xFD)
    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    component = safe_json(T2 / "component_ablation_summary.json", [])
    arch = safe_json(T2 / "architecture_summary.json", [])

    slides = []

    # 1 Cover
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF2, 0xF6, 0xF4)
    bg.line.fill.background()
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.28), SLIDE_H)
    band.fill.solid()
    band.fill.fore_color.rgb = GREEN
    band.line.fill.background()
    add_text(slide, "人工智能综合实训 II", Inches(0.95), Inches(1.20), Inches(8.8), Inches(0.55),
             size=24, color=GREEN_DARK, bold=True)
    add_text(slide, "机器学习项目答辩汇报", Inches(0.95), Inches(1.82), Inches(10.2), Inches(0.75),
             size=34, color=INK, bold=True)
    add_text(slide, "任务一：白葡萄酒质量分类  ·  任务二：手写 AlexNet 服饰图像分类",
             Inches(0.98), Inches(2.70), Inches(10.5), Inches(0.35), size=15, color=MUTED)
    add_card(slide, Inches(0.98), Inches(4.25), Inches(3.05), Inches(1.00),
             "任务一最佳模型", "随机森林\nmacro-F1 0.734", accent=GREEN)
    add_card(slide, Inches(4.35), Inches(4.25), Inches(3.05), Inches(1.00),
             "任务二主模型", "AlexNet\nAcc 94.38%", accent=BLUE)
    add_card(slide, Inches(7.72), Inches(4.25), Inches(3.05), Inches(1.00),
             "实验输出", "两篇论文 + 源码\n图表与复现说明", accent=AMBER, body_size=15)
    add_text(slide, "小组成员：雷正、蔡铭飞、冼嘉谦    指导教师：赵静    2026 年 6 月",
             Inches(0.98), Inches(6.45), Inches(11.0), Inches(0.30), size=12.5, color=MUTED)
    slides.append(slide)

    # 2 Agenda
    slide = blank_slide(prs)
    add_title(slide, "汇报目录", "围绕问题定义、方法实现、实验结果和误差分析展开", "Overview")
    agenda = [
        ("01", "项目整体设计", "两个任务分别覆盖传统机器学习与深度卷积网络"),
        ("02", "任务一：结构化数据分类", "白葡萄酒质量三分类、模型比较与有序误差分析"),
        ("03", "任务二：图像分类", "逐层手写 AlexNet、训练策略、消融与可解释性"),
        ("04", "结论与答辩要点", "完成度、不足、可解释回答与后续改进"),
    ]
    for i, (num, title, desc) in enumerate(agenda):
        y = Inches(1.55 + i * 1.15)
        add_card(slide, Inches(1.0), y, Inches(1.05), Inches(0.72), num, None, accent=GREEN)
        add_text(slide, num, Inches(1.05), y + Inches(0.12), Inches(0.95), Inches(0.35),
                 size=19, color=GREEN_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, Inches(2.30), y + Inches(0.03), Inches(4.0), Inches(0.30),
                 size=18, color=INK, bold=True)
        add_text(slide, desc, Inches(2.30), y + Inches(0.42), Inches(8.2), Inches(0.25),
                 size=12, color=MUTED)
    add_footer(slide, 2)
    slides.append(slide)

    # 3 Overall design
    slide = blank_slide(prs)
    add_title(slide, "项目整体完成情况", "两条技术路线均完成数据、模型、评估、分析和论文输出", "Scope")
    add_table(slide, Inches(0.75), Inches(1.55), Inches(11.85), Inches(1.55),
              ["任务", "数据集", "核心方法", "最终输出"],
              [
                  ["任务一", "UCI 白葡萄酒质量\n4898 样本，11 特征", "SVM / 决策树 / 随机森林 / 逻辑回归\n五折交叉验证 + 网格搜索", "独立论文、模型对比、混淆矩阵、特征重要性"],
                  ["任务二", "Fashion-MNIST\n70000 张灰度图，10 类", "PyTorch 基础算子逐层实现 AlexNet\nBN、增强、消融、Grad-CAM", "独立论文、训练曲线、消融实验、复杂度分析"],
              ],
              font_size=10.5)
    add_card(slide, Inches(0.85), Inches(3.75), Inches(3.55), Inches(1.35),
             "共同流程", "加载 → 查看 → 预处理 → 建模 → 评估 → 结论", accent=GREEN, body_size=16)
    add_card(slide, Inches(4.90), Inches(3.75), Inches(3.55), Inches(1.35),
             "评价侧重点", "不仅报告准确率，还分析类别混淆、误差结构和计算代价", accent=BLUE, body_size=15)
    add_card(slide, Inches(8.95), Inches(3.75), Inches(3.55), Inches(1.35),
             "复现能力", "脚本化训练、评估、绘图和论文/PPT 生成，输出路径清晰", accent=AMBER, body_size=15)
    add_footer(slide, 3)
    slides.append(slide)

    # 4 Method roadmap
    slide = blank_slide(prs)
    add_title(slide, "整体技术路线", "两个任务对应不同数据形态，但实验逻辑保持一致", "Pipeline")
    add_image_fit(slide, REPORT / "fig_ml_pipeline.png", Inches(0.65), Inches(1.55),
                  Inches(5.95), Inches(4.75))
    add_image_fit(slide, REPORT / "fig_train_pipeline.png", Inches(6.75), Inches(1.55),
                  Inches(5.95), Inches(4.75))
    add_text(slide, "任务一：结构化表格数据，重点在标准化、交叉验证和模型比较",
             Inches(0.80), Inches(6.40), Inches(5.4), Inches(0.35), size=11.5, color=MUTED)
    add_text(slide, "任务二：图像数据，重点在网络结构、训练策略和可解释性分析",
             Inches(6.90), Inches(6.40), Inches(5.4), Inches(0.35), size=11.5, color=MUTED)
    add_footer(slide, 4)
    slides.append(slide)

    # 5 Section task 1
    slide = blank_slide(prs)
    add_title(slide, "任务一：白葡萄酒质量三分类", "从理化指标预测低、中、高质量等级", "Task 1")
    add_card(slide, Inches(0.75), Inches(1.40), Inches(3.1), Inches(1.1),
             "数据规模", "4898 × 11", accent=GREEN)
    add_card(slide, Inches(4.05), Inches(1.40), Inches(3.1), Inches(1.1),
             "标签构造", "≤5 / =6 / ≥7", accent=BLUE)
    add_card(slide, Inches(7.35), Inches(1.40), Inches(3.1), Inches(1.1),
             "选择目标", "macro-F1", accent=AMBER)
    add_bullets(slide, [
        "原始质量评分集中在 5、6、7 分，极端评分样本较少。",
        "三分类能降低少数分值过稀带来的训练和评价波动。",
        "质量等级存在自然顺序，因此后续引入二次加权 Kappa、有序 MAE 和严重跨级误判数。",
    ], Inches(0.85), Inches(3.10), Inches(5.75), Inches(1.45), size=14.5)
    add_image_fit(slide, T1 / "eda_quality_dist.png", Inches(7.05), Inches(2.85),
                  Inches(5.15), Inches(2.15))
    add_image_fit(slide, T1 / "eda_corr_heatmap.png", Inches(0.85), Inches(4.95),
                  Inches(5.75), Inches(1.65))
    add_text(slide, "相关性热力图提示理化指标之间存在联合影响，适合比较非线性模型与集成模型。",
             Inches(7.05), Inches(5.35), Inches(5.0), Inches(0.55), size=12.5, color=MUTED)
    add_footer(slide, 5)
    slides.append(slide)

    # 6 Task1 modeling design
    slide = blank_slide(prs)
    add_title(slide, "任务一建模设计", "重点控制数据泄漏，并用宏平均 F1 做模型选择", "Design")
    add_bullets(slide, [
        "预处理：StandardScaler 封装在 Pipeline 中，只在训练折内拟合。",
        "模型比较：SVM、决策树、随机森林、逻辑回归。",
        "调参方式：五折交叉验证 + 网格搜索，scoring=f1_macro。",
        "测试评价：准确率、宏 P/R/F1、平衡准确率、Kappa、MCC、ROC-AUC。",
    ], Inches(0.85), Inches(1.55), Inches(5.55), Inches(2.10), size=15)
    add_table(slide, Inches(0.85), Inches(4.00), Inches(11.7), Inches(1.55),
              ["模型", "搜索空间", "真实最优参数"],
              [
                  ["SVM", "C∈{1,10}; gamma∈{scale,0.1}", "C=10; gamma=0.1"],
                  ["决策树", "深度∈{None,10,20}; 叶子∈{1,5}", "深度 None; 叶子 1"],
                  ["随机森林", "树数∈{200,400}; 深度∈{None,20}", "树数 200; 深度 20"],
                  ["逻辑回归", "C∈{0.5,1,10}", "C=0.5"],
              ],
              font_size=9.5)
    add_card(slide, Inches(7.05), Inches(1.55), Inches(5.1), Inches(1.45),
             "答辩重点", "参数表已与真实运行输出同步；指标与 outputs/metrics.csv 保持一致。", accent=RED, body_size=14)
    add_footer(slide, 6)
    slides.append(slide)

    # 7 Task1 results
    slide = blank_slide(prs)
    add_title(slide, "任务一实验结果：随机森林表现最好", "集成模型更适合理化指标的非线性组合关系", "Result")
    add_image_fit(slide, T1 / "model_compare.png", Inches(0.70), Inches(1.45),
                  Inches(6.0), Inches(4.45))
    add_table(slide, Inches(7.05), Inches(1.45), Inches(5.35), Inches(2.15),
              ["模型", "Acc", "macro-F1"],
              [
                  ["SVM", "0.617", "0.622"],
                  ["决策树", "0.650", "0.645"],
                  ["随机森林", "0.735", "0.734"],
                  ["逻辑回归", "0.522", "0.520"],
              ],
              font_size=11)
    add_bullets(slide, [
        "逻辑回归较弱，说明标签边界并非简单线性可分。",
        "单棵决策树有提升，但对训练样本扰动更敏感。",
        "随机森林通过多树投票降低方差，在基础指标和进阶指标上均占优。",
    ], Inches(7.10), Inches(4.05), Inches(5.15), Inches(1.55), size=13.5)
    add_footer(slide, 7)
    slides.append(slide)

    # 8 Task1 error and interpretation
    slide = blank_slide(prs)
    add_title(slide, "任务一误差分析与特征解释", "错误主要集中在相邻质量等级之间", "Analysis")
    add_image_fit(slide, T1 / "cm_随机森林.png", Inches(0.70), Inches(1.45),
                  Inches(4.15), Inches(3.35))
    add_image_fit(slide, T1 / "rf_feature_importance.png", Inches(5.10), Inches(1.45),
                  Inches(4.15), Inches(3.35))
    add_image_fit(slide, T1 / "ordinal_improvement.png", Inches(9.50), Inches(1.45),
                  Inches(3.05), Inches(3.35))
    add_bullets(slide, [
        "混淆矩阵：主要误差出现在“差-中”“中-好”之间，严重跨级误判较少。",
        "特征重要性：酒精度、密度、挥发性酸度等特征对树模型划分贡献较大。",
        "有序建模：回归再分级使 macro-F1 从 0.7336 提升至 0.7533，严重误判从 7 降至 4。",
    ], Inches(0.90), Inches(5.35), Inches(11.5), Inches(0.95), size=12.7)
    add_footer(slide, 8)
    slides.append(slide)

    # 9 Section task2
    slide = blank_slide(prs)
    add_title(slide, "任务二：基于 AlexNet 的服饰图像分类", "逐层手写 CNN，不调用预置 AlexNet 模型", "Task 2")
    add_card(slide, Inches(0.75), Inches(1.40), Inches(3.1), Inches(1.05),
             "数据集", "Fashion-MNIST\n10 类服饰", accent=GREEN, body_size=16)
    add_card(slide, Inches(4.05), Inches(1.40), Inches(3.1), Inches(1.05),
             "输入适配", "1×28×28 → 1×224×224", accent=BLUE, body_size=15)
    add_card(slide, Inches(7.35), Inches(1.40), Inches(3.1), Inches(1.05),
             "主模型结果", "Acc 94.38%\nmacro-F1 0.9437", accent=AMBER, body_size=15)
    add_image_fit(slide, T2 / "samples.png", Inches(0.85), Inches(3.05),
                  Inches(5.55), Inches(2.75))
    add_bullets(slide, [
        "输入通道从 RGB 三通道改为单通道灰度。",
        "最后一层输出从 1000 类改为 10 类。",
        "保留五个卷积层和三层全连接层主体结构，用 BatchNorm 作为主模型配置。",
    ], Inches(6.85), Inches(3.20), Inches(5.3), Inches(1.60), size=14.5)
    add_footer(slide, 9)
    slides.append(slide)

    # 10 AlexNet structure
    slide = blank_slide(prs)
    add_title(slide, "手写 AlexNet 结构与训练策略", "经典结构适配小尺寸灰度图像", "Model")
    add_image_fit(slide, REPORT / "fig_alexnet_arch.png", Inches(0.70), Inches(1.45),
                  Inches(6.1), Inches(4.15))
    add_table(slide, Inches(7.05), Inches(1.45), Inches(5.25), Inches(2.10),
              ["模块", "结构", "输出"],
              [
                  ["Conv1", "11×11 Conv + BN + Pool", "96×27×27"],
                  ["Conv2", "5×5 Conv + BN + Pool", "256×13×13"],
                  ["Conv3-5", "3×3 Conv + BN/ReLU", "256×6×6"],
                  ["Classifier", "Dropout + FC4096 + FC10", "10"],
              ],
              font_size=9.5)
    add_bullets(slide, [
        "优化器：SGD，momentum=0.9，weight decay=5e-4。",
        "训练：40 轮，batch size=128，cosine annealing。",
        "正则化：BatchNorm、Dropout、随机翻转和小角度旋转。",
    ], Inches(7.10), Inches(3.95), Inches(5.10), Inches(1.45), size=13.5)
    add_footer(slide, 10)
    slides.append(slide)

    # 11 Task2 main results
    slide = blank_slide(prs)
    add_title(slide, "任务二主模型测试结果", "总体性能较高，错误集中在视觉相似上装类别", "Result")
    add_image_fit(slide, T2 / "curves.png", Inches(0.65), Inches(1.45),
                  Inches(5.55), Inches(3.30))
    add_image_fit(slide, T2 / "confusion_matrix.png", Inches(6.45), Inches(1.45),
                  Inches(5.55), Inches(3.30))
    add_table(slide, Inches(0.90), Inches(5.25), Inches(5.35), Inches(0.80),
              ["Acc", "Precision", "Recall", "macro-F1", "ROC-AUC", "Kappa"],
              [["0.9438", "0.9437", "0.9438", "0.9437", "0.9972", "0.9376"]],
              font_size=9.5)
    add_bullets(slide, [
        "裤子、包、鞋类整体轮廓差异大，分类效果较好。",
        "T 恤、套衫、外套、衬衫在低分辨率灰度图中局部差异弱，容易混淆。",
    ], Inches(6.55), Inches(5.25), Inches(5.4), Inches(0.8), size=12.2)
    add_footer(slide, 11)
    slides.append(slide)

    # 12 Ablation
    slide = blank_slide(prs)
    add_title(slide, "任务二消融实验", "充分训练和 BatchNorm 是主要性能增益来源", "Ablation")
    ab_rows = [
        ["15轮 LRN", "0.9155", "0.9153", "0.7700"],
        ["40轮 LRN", "0.9314", "0.9311", "0.7926"],
        ["40轮 BN", "0.9432", "0.9429", "0.8262"],
        ["BN+增强", "0.9438", "0.9437", "0.8281"],
    ]
    if component:
        ab_rows = [[r["name"].replace("最小配置（15轮，LRN，无增强）", "15轮 LRN")
                    .replace("+ 充分训练（40轮，余弦退火）", "40轮 LRN")
                    .replace("+ BatchNorm 替代 LRN", "40轮 BN")
                    .replace("+ 数据增强（主模型）", "BN+增强"),
                    f"{r['acc']:.4f}", f"{r['f1']:.4f}", f"{r['shirt_f1']:.4f}"]
                   for r in component]
    add_table(slide, Inches(0.80), Inches(1.45), Inches(5.9), Inches(1.65),
              ["配置", "Acc", "macro-F1", "衬衫 F1"], ab_rows, font_size=10.5)
    add_image_fit(slide, T2 / "improve_compare.png", Inches(7.00), Inches(1.25),
                  Inches(5.25), Inches(3.95))
    add_bullets(slide, [
        "15 轮训练不足以充分发挥 AlexNet 容量。",
        "BatchNorm 替代 LRN 后准确率提升明显，同时困难类 F1 改善。",
        "数据增强在当前设置下提升较小，更像轻度正则化。",
    ], Inches(0.95), Inches(4.05), Inches(5.55), Inches(1.30), size=13.5)
    add_footer(slide, 12)
    slides.append(slide)

    # 13 Complexity
    slide = blank_slide(prs)
    add_title(slide, "复杂度与架构对照", "AlexNet 能取得较好性能，但在小图像任务上计算冗余明显", "Efficiency")
    add_table(slide, Inches(0.75), Inches(1.45), Inches(6.15), Inches(1.40),
              ["模型", "参数量", "FLOPs", "CPU 延迟", "Acc"],
              [
                  ["AlexNet@224", "58.302M", "1063.5M", "15.437ms", "94.38%"],
                  ["SimpleCNN@28", "0.391M", "7.9M", "1.496ms", "约91.8%"],
                  ["倍数", "149.1×", "134.6×", "10.3×", "-"],
              ],
              font_size=9.5)
    add_image_fit(slide, T2 / "arch_compare.png", Inches(7.15), Inches(1.40),
                  Inches(5.05), Inches(3.35))
    arch_text = "ResNetSmall@224：Acc 94.92%，参数量 2.777M"
    if arch:
        resnet = next((r for r in arch if "ResNet" in r["name"]), None)
        if resnet:
            arch_text = f"ResNetSmall@224：Acc {resnet['acc']:.4f}，参数量 {resnet['params_M']:.3f}M"
    add_card(slide, Inches(0.90), Inches(3.65), Inches(5.75), Inches(1.12),
             "架构对照结论", arch_text, accent=BLUE, body_size=15)
    add_bullets(slide, [
        "Fashion-MNIST 原图只有 28×28，Resize 到 224×224 是结构适配，不会增加真实细节。",
        "全连接分类器贡献大量参数，适合教学理解经典结构，但不一定是部署最优方案。",
    ], Inches(7.20), Inches(5.05), Inches(4.9), Inches(0.95), size=12.8)
    add_footer(slide, 13)
    slides.append(slide)

    # 14 Explainability
    slide = blank_slide(prs)
    add_title(slide, "Focal Loss 与 Grad-CAM 误差分析", "上装混淆更接近细粒度视觉差异不足，而非类别数量不均衡", "Explain")
    add_table(slide, Inches(0.80), Inches(1.45), Inches(5.05), Inches(1.15),
              ["损失函数", "Acc", "macro-F1", "衬衫 F1", "外套 F1"],
              [
                  ["交叉熵", "0.9247", "0.9245", "0.7709", "0.8812"],
                  ["Focal Loss", "0.9167", "0.9163", "0.7491", "0.8637"],
              ],
              font_size=10.2)
    add_image_fit(slide, T2 / "gradcam.png", Inches(6.25), Inches(1.35),
                  Inches(5.85), Inches(3.95))
    add_bullets(slide, [
        "Focal Loss 未提升，说明错误来源不主要是类别不均衡。",
        "Grad-CAM 显示模型主要关注服饰主体轮廓。",
        "T 恤、套衫、外套、衬衫依赖领口、袖口、开襟等局部细节，低分辨率下信息不足。",
    ], Inches(0.90), Inches(3.45), Inches(5.10), Inches(1.75), size=13.2)
    add_footer(slide, 14)
    slides.append(slide)

    # 15 Conclusions
    slide = blank_slide(prs)
    add_title(slide, "总结与不足", "两个项目均完成了从实验到论文的闭环", "Conclusion")
    add_card(slide, Inches(0.80), Inches(1.35), Inches(3.7), Inches(1.30),
             "任务一结论", "随机森林最优；有序回归分级可减少严重跨级误判。", accent=GREEN, body_size=14)
    add_card(slide, Inches(4.85), Inches(1.35), Inches(3.7), Inches(1.30),
             "任务二结论", "手写 AlexNet 达到 94.38%；上装类别仍是主要难点。", accent=BLUE, body_size=14)
    add_card(slide, Inches(8.90), Inches(1.35), Inches(3.7), Inches(1.30),
             "工程产出", "训练、评估、绘图、报告和 PPT 均脚本化生成。", accent=AMBER, body_size=14)
    add_bullets(slide, [
        "不足 1：部分对照实验仍以单随机种子为主，后续可报告均值和标准差。",
        "不足 2：任务一可加入 XGBoost、LightGBM、专门有序分类模型。",
        "不足 3：任务二可比较更适合小图像的轻量 CNN、ResNet 变体或注意力模块。",
    ], Inches(1.05), Inches(3.35), Inches(10.9), Inches(1.40), size=14.2)
    add_text(slide, "核心答辩口径：我们不只展示最高指标，也分析了为什么会错、错误是否严重、以及模型是否高效。",
             Inches(1.00), Inches(5.70), Inches(11.2), Inches(0.38), size=15, color=GREEN_DARK, bold=True,
             align=PP_ALIGN.CENTER)
    add_footer(slide, 15)
    slides.append(slide)

    # 16 Q&A preparation
    slide = blank_slide(prs)
    add_title(slide, "答辩追问准备", "建议重点记住这些解释口径", "Q&A")
    qa_rows = [
        ["为什么任务一随机森林最好？", "表格数据存在非线性特征组合，多树集成降低单树方差。"],
        ["为什么做有序回归对照？", "葡萄酒质量有低中高顺序，相邻误判和跨级误判含义不同。"],
        ["为什么 AlexNet 有冗余？", "28×28 图像放大到 224×224 只适配结构，不增加真实细节；全连接层参数量大。"],
        ["为什么上装类别容易错？", "低分辨率灰度图缺少颜色和纹理，T恤/套衫/外套/衬衫局部差异弱。"],
        ["为什么 BatchNorm 有提升？", "稳定中间层分布，使大模型在 40 轮训练中更充分、更平稳地收敛。"],
    ]
    add_table(slide, Inches(0.75), Inches(1.40), Inches(11.85), Inches(4.35),
              ["可能问题", "回答要点"], qa_rows, font_size=10.5)
    add_text(slide, "汇报结束，敬请批评指正", Inches(0.95), Inches(6.25),
             Inches(11.4), Inches(0.45), size=24, color=GREEN_DARK, bold=True,
             align=PP_ALIGN.CENTER)
    add_footer(slide, 16)
    slides.append(slide)

    prs.save(OUT)
    return OUT, len(slides)


if __name__ == "__main__":
    out, n_slides = build_deck()
    print(f"[PPT] generated: {out}")
    print(f"[PPT] slides: {n_slides}")
